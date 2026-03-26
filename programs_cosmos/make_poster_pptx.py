"""
Generate a PowerPoint poster for supernova detection presentation.

Usage:
    pip install python-pptx
    python make_poster_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Poster size: 48 x 36 inches (landscape A0-ish)
POSTER_W = Inches(48)
POSTER_H = Inches(36)

# Colors
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x2E, 0x50, 0x90)
LIGHT_BLUE = RGBColor(0xE8, 0xEE, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
RED = RGBColor(0xC0, 0x30, 0x30)
GREEN = RGBColor(0x20, 0x80, 0x40)
ORANGE = RGBColor(0xE0, 0x80, 0x20)
LIGHT_GREEN = RGBColor(0xE0, 0xF5, 0xE0)
LIGHT_YELLOW = RGBColor(0xFF, 0xF8, 0xE0)
LIGHT_ORANGE = RGBColor(0xFF, 0xF0, 0xE0)
LIGHT_RED = RGBColor(0xFF, 0xE8, 0xE8)
GRAY = RGBColor(0x60, 0x60, 0x60)


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height, lines, font_name="Calibri"):
    """Add a textbox with multiple lines, each with its own formatting.
    lines: list of (text, font_size, bold, color) tuples
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, font_size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(4)

    return txBox


def add_block(slide, left, top, width, height, title, title_size=28):
    """Add a colored block with title bar."""
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                left, top, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BLUE
    bg.line.color.rgb = BLUE
    bg.line.width = Pt(2)

    # Title bar
    title_h = Inches(0.7)
    tb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                left, top, width, title_h)
    tb.fill.solid()
    tb.fill.fore_color.rgb = BLUE
    tb.line.color.rgb = NAVY
    tb.line.width = Pt(2)

    # Title text
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.05),
                                     width - Inches(0.4), title_h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    return left + Inches(0.3), top + title_h + Inches(0.15)


def add_step_box(slide, left, top, width, height, text, fill_color):
    """Add a pipeline step box."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = BLUE
    box.line.width = Pt(1.5)

    txBox = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.05),
                                     width - Inches(0.2), height - Inches(0.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = False
    p.font.color.rgb = BLACK
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    return box


def add_arrow_down(slide, cx, top, length):
    """Add a downward arrow."""
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                   cx - Inches(0.15), top,
                                   Inches(0.3), length)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BLUE
    arrow.line.color.rgb = BLUE


def add_cnn_layer(slide, left, top, width, height, text, fill_color):
    """Add a CNN layer box."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = GRAY
    box.line.width = Pt(1)

    txBox = slide.shapes.add_textbox(left + Inches(0.05), top,
                                     width - Inches(0.1), height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = BLACK
    p.font.name = "Calibri"


def main():
    prs = Presentation()
    prs.slide_width = POSTER_W
    prs.slide_height = POSTER_H

    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # ================================================================
    # TITLE BAR
    # ================================================================
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       0, 0, POSTER_W, Inches(3.5))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()

    # Title
    add_textbox(slide, Inches(1), Inches(0.3), Inches(46), Inches(1.5),
                "Automated Detection of Supernovae in the COSMOS Field\n"
                "by Comparing JWST and HST Images with Deep Learning",
                font_size=52, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    # Author
    add_textbox(slide, Inches(1), Inches(2.2), Inches(46), Inches(1),
                "N. Suzuki",
                font_size=36, bold=False, color=RGBColor(0xCC, 0xDD, 0xFF),
                alignment=PP_ALIGN.CENTER)

    # ================================================================
    # BLOCK 1: Introduction (top-left)
    # ================================================================
    b1_left, b1_top = Inches(0.5), Inches(4.0)
    b1_w, b1_h = Inches(15.3), Inches(9.5)
    cx, cy = add_block(slide, b1_left, b1_top, b1_w, b1_h,
                        "1. Introduction")

    add_rich_textbox(slide, cx, cy, b1_w - Inches(0.6), b1_h - Inches(1.2), [
        ("Supernovae (SNe) are rare transient events — bright stellar "
         "explosions that appear as new point sources near their host galaxies.",
         18, False, BLACK),
        ("", 10, False, BLACK),
        ("Challenge:", 20, True, NAVY),
        ("• The COSMOS field was observed by both HST (2004–2005) "
         "and JWST (2023–2024), spanning ~20 years", 18, False, BLACK),
        ("• Images from different telescopes with different pixel scales, "
         "filters, and resolutions", 18, False, BLACK),
        ("• Manual inspection of ~212,000 objects is infeasible", 18, False, BLACK),
        ("• SN rate is very low: ~1 per 1,000 objects", 18, False, RED),
        ("", 10, False, BLACK),
        ("Goal:", 20, True, NAVY),
        ("Build an automated CNN-based pipeline to detect SN candidates "
         "by comparing HST and JWST image triplets, validated against "
         "23 visually confirmed supernovae.", 18, False, BLACK),
    ])

    # ================================================================
    # BLOCK 2: Data (top-right)
    # ================================================================
    b2_left, b2_top = Inches(16.3), Inches(4.0)
    b2_w, b2_h = Inches(15.3), Inches(9.5)
    cx, cy = add_block(slide, b2_left, b2_top, b2_w, b2_h,
                        "2. Data")

    add_rich_textbox(slide, cx, cy, b2_w - Inches(0.6), b2_h - Inches(1.2), [
        ("HST COSMOS-ACS:  212,460 PNG cutouts", 18, True, BLACK),
        ("    Filename: {ID}_{Field}_hstcosmos.png", 16, False, GRAY),
        ("", 8, False, BLACK),
        ("JWST COSMOS-Web:  424,941 PNG cutouts (2 filters/object)", 18, True, BLACK),
        ("    Filenames: {ID}_{Field}_jwst1.png, {ID}_{Field}_jwst2.png", 16, False, GRAY),
        ("", 8, False, BLACK),
        ("Matching: By object ID (field names may differ)", 18, False, BLACK),
        ("", 8, False, BLACK),
        ("Image properties:", 20, True, NAVY),
        ("• Galaxy centered in each PNG cutout", 18, False, BLACK),
        ("• Grayscale, resized to 64×64 pixels", 18, False, BLACK),
        ("• Central 50% cropped to focus on galaxy", 18, False, BLACK),
        ("", 8, False, BLACK),
        ("Training/validation: 23 confirmed supernovae from visual inspection",
         18, False, RED),
        ("", 8, False, BLACK),
        ("Dataset          Count        Telescope", 16, True, NAVY),
        ("HST images       212,460     HST ACS/WFC", 16, False, BLACK),
        ("JWST images     424,941     JWST NIRCam", 16, False, BLACK),
        ("Matched triplets  ~212,000   HST + JWST1 + JWST2", 16, False, BLACK),
        ("Known SNe         23            Visual inspection", 16, False, BLACK),
        ("Expected rate     ~1/1000", 16, False, RED),
    ])

    # ================================================================
    # BLOCK 3: Physical Constraints (top-far-right)
    # ================================================================
    b3_left, b3_top = Inches(32.1), Inches(4.0)
    b3_w, b3_h = Inches(15.3), Inches(9.5)
    cx, cy = add_block(slide, b3_left, b3_top, b3_w, b3_h,
                        "3. Physical Constraints")

    add_rich_textbox(slide, cx, cy, b3_w - Inches(0.6), Inches(2.5), [
        ("(a) SN appears near host galaxy", 20, True, NAVY),
        ("A supernova is an exploding star within a galaxy. It must appear "
         "near the center of each PNG (where the host galaxy is). "
         "We crop the central 50% and keep sources within 25% of center.",
         18, False, BLACK),
        ("", 10, False, BLACK),
        ("(b) SN must be consistent across filters:", 20, True, NAVY),
    ])

    # JWST SN diagram
    diag_top = cy + Inches(3.0)
    diag_left = cx + Inches(0.5)

    # JWST SN boxes
    add_cnn_layer(slide, diag_left, diag_top, Inches(3.5), Inches(0.7),
                  "HST: No SN", LIGHT_GREEN)
    add_cnn_layer(slide, diag_left + Inches(5), diag_top - Inches(0.4),
                  Inches(3.5), Inches(0.7), "JWST1: SN present", LIGHT_RED)
    add_cnn_layer(slide, diag_left + Inches(5), diag_top + Inches(0.5),
                  Inches(3.5), Inches(0.7), "JWST2: SN present", LIGHT_RED)
    add_textbox(slide, diag_left + Inches(9.5), diag_top,
                Inches(4), Inches(0.7),
                "→ JWST Supernova", font_size=18, bold=True, color=RED)

    # HST SN boxes
    ht = diag_top + Inches(2.0)
    add_cnn_layer(slide, diag_left, ht, Inches(3.5), Inches(0.7),
                  "HST: SN present", LIGHT_RED)
    add_cnn_layer(slide, diag_left + Inches(5), ht - Inches(0.4),
                  Inches(3.5), Inches(0.7), "JWST1: No SN", LIGHT_GREEN)
    add_cnn_layer(slide, diag_left + Inches(5), ht + Inches(0.5),
                  Inches(3.5), Inches(0.7), "JWST2: No SN", LIGHT_GREEN)
    add_textbox(slide, diag_left + Inches(9.5), ht,
                Inches(4), Inches(0.7),
                "→ HST Supernova", font_size=18, bold=True, color=RED)

    # Bullets below diagram
    add_rich_textbox(slide, cx, ht + Inches(1.5),
                     b3_w - Inches(0.6), Inches(2.5), [
        ("• JWST SN: new source in both JWST1 and JWST2 "
         "(same position), absent in HST", 18, False, BLACK),
        ("• HST SN: source in HST only, absent in both "
         "JWST1 and JWST2 (SN faded)", 18, False, BLACK),
        ("• Requiring both JWST filters eliminates single-filter "
         "artifacts and noise", 18, False, GREEN),
    ])

    # ================================================================
    # BLOCK 4: CNN Architecture (middle-left)
    # ================================================================
    b4_left, b4_top = Inches(0.5), Inches(14.0)
    b4_w, b4_h = Inches(15.3), Inches(11.5)
    cx, cy = add_block(slide, b4_left, b4_top, b4_w, b4_h,
                        "4. CNN Architecture & Training")

    # Input
    add_textbox(slide, cx, cy, Inches(5), Inches(0.6),
                "Input: 3-channel 64×64 (HST + JWST1 + JWST2)",
                font_size=18, bold=True, color=NAVY)
    add_textbox(slide, cx, cy + Inches(0.5), Inches(5), Inches(0.6),
                "Output: P(supernova) ∈ [0, 1]",
                font_size=18, bold=True, color=NAVY)

    # CNN layers diagram
    layer_left = cx + Inches(0.5)
    layer_w = Inches(6)
    layer_h = Inches(0.6)
    gap = Inches(0.85)
    arrow_h = Inches(0.2)
    y = cy + Inches(1.5)

    layers = [
        ("Input: 3 × 64 × 64", LIGHT_GREEN),
        ("Conv2D(16) + BN + ReLU → MaxPool → 32×32", LIGHT_YELLOW),
        ("Conv2D(32) + BN + ReLU → MaxPool → 16×16", LIGHT_YELLOW),
        ("Conv2D(64) + BN + ReLU → MaxPool → 8×8", LIGHT_YELLOW),
        ("Conv2D(128) + BN + ReLU → AvgPool → 1×1", LIGHT_YELLOW),
        ("FC(128→64) + ReLU + Dropout(0.3) + FC(64→1)", LIGHT_ORANGE),
        ("Sigmoid → P(SN)", LIGHT_RED),
    ]

    for i, (text, color) in enumerate(layers):
        add_cnn_layer(slide, layer_left, y, layer_w, layer_h, text, color)
        if i < len(layers) - 1:
            add_arrow_down(slide, layer_left + layer_w / 2,
                          y + layer_h + Inches(0.02), arrow_h)
        y += gap

    # Training strategy
    strat_left = cx + Inches(7.5)
    add_rich_textbox(slide, strat_left, cy + Inches(1.2),
                     Inches(6.5), Inches(8), [
        ("Training Strategy:", 20, True, NAVY),
        ("", 8, False, BLACK),
        ("• 23 known SN oversampled ×50", 17, False, BLACK),
        ("  as real positive examples", 17, False, BLACK),
        ("", 6, False, BLACK),
        ("• 5,000 artificial SN injected as", 17, False, BLACK),
        ("  Gaussian sources near galaxy center", 17, False, BLACK),
        ("  (both JWST-type and HST-type)", 17, False, BLACK),
        ("", 6, False, BLACK),
        ("• 100,000 negatives (20:1 ratio)", 17, False, BLACK),
        ("  from non-SN objects", 17, False, BLACK),
        ("", 6, False, BLACK),
        ("• Loss: BCEWithLogits", 17, False, BLACK),
        ("• Optimizer: Adam, LR = 10⁻³", 17, False, BLACK),
        ("• LR scheduler: ReduceLROnPlateau", 17, False, BLACK),
        ("• Augmentation: flips + 90° rotations", 17, False, BLACK),
        ("• 50 epochs, best model by", 17, False, BLACK),
        ("  known SN recovery rate", 17, False, BLACK),
    ])

    # ================================================================
    # BLOCK 5: Detection Pipeline (middle-center)
    # ================================================================
    b5_left, b5_top = Inches(16.3), Inches(14.0)
    b5_w, b5_h = Inches(15.3), Inches(11.5)
    cx, cy = add_block(slide, b5_left, b5_top, b5_w, b5_h,
                        "5. Detection Pipeline")

    step_w = Inches(12)
    step_h = Inches(1.0)
    step_gap = Inches(1.5)
    step_left = cx + Inches(1.2)
    y = cy + Inches(0.5)

    steps = [
        ("Step 1:  Match HST + JWST1 + JWST2 triplets by object ID", LIGHT_GREEN),
        ("Step 2:  Train CNN on known SN + artificial SN", LIGHT_YELLOW),
        ("Step 3:  Validate — recover ALL 23 known supernovae", LIGHT_ORANGE),
        ("Step 4:  Apply to all ~212,000 objects (write results immediately)",
         LIGHT_RED),
        ("Step 5:  Verify known SN in results; check rate ~1/1000",
         LIGHT_BLUE),
    ]

    for i, (text, color) in enumerate(steps):
        add_step_box(slide, step_left, y, step_w, step_h, text, color)
        if i < len(steps) - 1:
            add_arrow_down(slide, step_left + step_w / 2,
                          y + step_h + Inches(0.05), Inches(0.4))
        y += step_gap

    add_rich_textbox(slide, cx, y + Inches(0.3),
                     b5_w - Inches(0.6), Inches(1.5), [
        ("Candidates are written to CSV immediately as they are found,",
         17, False, BLACK),
        ("allowing real-time monitoring during the multi-hour run.",
         17, False, BLACK),
    ])

    # ================================================================
    # BLOCK 6: Threshold & Artificial SN (middle-right)
    # ================================================================
    b6_left, b6_top = Inches(32.1), Inches(14.0)
    b6_w, b6_h = Inches(15.3), Inches(11.5)
    cx, cy = add_block(slide, b6_left, b6_top, b6_w, b6_h,
                        "6. Threshold Selection & Artificial SN")

    add_rich_textbox(slide, cx, cy, b6_w - Inches(0.6), b6_h - Inches(1.0), [
        ("Automatic threshold selection:", 20, True, NAVY),
        ("The threshold is set to 90% of the lowest probability among "
         "the 23 known SNe. This guarantees 100% recovery of all "
         "confirmed SN.", 18, False, BLACK),
        ("", 8, False, BLACK),
        ("    θ = 0.9 × min P(SN) over known SNe", 20, True, RED),
        ("", 10, False, BLACK),
        ("Artificial supernova injection:", 20, True, NAVY),
        ("Gaussian point sources injected near the galaxy center "
         "to augment the small training set:", 18, False, BLACK),
        ("", 6, False, BLACK),
        ("    I_SN(x,y) = I_orig(x,y) + A·exp(−((x−x₀)²+(y−y₀)²) / 2σ²)",
         17, False, BLUE),
        ("", 8, False, BLACK),
        ("Parameters:", 18, True, NAVY),
        ("• (x₀, y₀): random position within 25% of image center",
         17, False, BLACK),
        ("• A ∈ [0.3, 0.9]: peak brightness (random)", 17, False, BLACK),
        ("• σ ∈ [1.0, 3.0] pixels: PSF width (random)", 17, False, BLACK),
        ("", 6, False, BLACK),
        ("For JWST SN:", 18, True, NAVY),
        ("• Injected at same position in both JWST filters", 17, False, BLACK),
        ("• Slightly different brightness: A₂ = A₁ × U[0.7, 1.3]", 17, False, BLACK),
        ("", 10, False, BLACK),
        ("[Space for artificial SN example images]", 18, True, GRAY),
    ])

    # ================================================================
    # BLOCK 7: Results placeholder (bottom)
    # ================================================================
    b7_left, b7_top = Inches(0.5), Inches(26.0)
    b7_w, b7_h = Inches(46.9), Inches(9.5)
    cx, cy = add_block(slide, b7_left, b7_top, b7_w, b7_h,
                        "7. Results")

    add_textbox(slide, cx + Inches(12), cy + Inches(3), Inches(20), Inches(2),
                "[Second half — to be filled with result images,\n"
                "candidate gallery, detection statistics, and recovery plots]",
                font_size=28, bold=False, color=GRAY,
                alignment=PP_ALIGN.CENTER)

    # Save
    output_path = "poster_supernova.pptx"
    prs.save(output_path)
    print(f"Poster saved to {output_path}")
    print(f"Size: 48 × 36 inches (standard poster size)")
    print(f"\nOpen in PowerPoint to edit and add result images.")


if __name__ == "__main__":
    main()
